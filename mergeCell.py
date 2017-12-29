from pptx import Presentation
from pptx.util import Inches


# copy from: https://groups.google.com/forum/#!topic/python-pptx/cVRP9sSpEjA

# merge cells vertically
def mergeCellsVertically(table, start_row_idx, end_row_idx, col_idx):
    row_count = end_row_idx - start_row_idx + 1
    column_cells = [r.cells[col_idx] for r in table.rows][start_row_idx:]

    column_cells[0]._tc.set('rowSpan', str(row_count))
    for c in column_cells[1:]:
        c._tc.set('vMerge', '1')


# merge cells horizontally
def mergeCellsHorizontally(table, row_idx, start_col_idx, end_col_idx):
    col_count = end_col_idx - start_col_idx + 1
    row_cells = [c for c in table.rows[row_idx].cells][start_col_idx:end_col_idx]
    row_cells[0]._tc.set('gridSpan', str(col_count))
    for c in row_cells[1:]:
        c._tc.set('hMerge', '1')


# the workaround function to merge cells in a table
def mergeCells(table, start_row_idx, end_row_idx, start_col_idx, end_col_idx):
    for col_idx in range(start_col_idx,end_col_idx+1):
        mergeCellsVertically(table,start_row_idx, end_row_idx, col_idx)
    for row_idx in range(start_row_idx,end_row_idx+1):
        mergeCellsHorizontally(table, row_idx, start_col_idx, end_col_idx)


if __name__ == '__main__':
    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = 'Adding a Table'

    rows = cols = 4
    top = Inches(2.0)
    left = Inches(1.0)
    width = Inches(8.0)
    height = Inches(1)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.0)

    # write column headings
    table.cell(0, 0).text = 'Column1'
    table.cell(0, 1).text = 'Column2'
    table.cell(0, 2).text = 'Column3'
    table.cell(0, 3).text = 'Column4'

    # write body cells
    table.cell(1, 0).text = 'Merged vertically and horizontally'
    table.cell(1, 2).text = 'Cell 1-3'
    table.cell(2, 2).text = 'Cell 2-3'
    table.cell(3, 2).text = 'Cell 3-3'
    table.cell(1, 3).text = 'Merged vertically'
    table.cell(3, 0).text = 'Merged horizontally'
    table.cell(3, 3).text = 'Cell 3-4'

    # merge cells vertically
    mergeCellsVertically(table=table, start_row_idx=1, end_row_idx=2, col_idx=3)
    # merge cells horizontally
    mergeCellsHorizontally(table=table, row_idx=3, start_col_idx=0, end_col_idx=1)
    # merge cells vertically and horizontally
    mergeCells(table=table, start_row_idx=1, end_row_idx=2, start_col_idx=0, end_col_idx=1)
    prs.save('test.pptx')
