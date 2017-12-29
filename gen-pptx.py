import json
import os.path
from pprint import pprint

import chardet
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from operateSlide import move_slide, duplicate_slide
from mergeCell import mergeCellsVertically, mergeCellsHorizontally, mergeCells


def setup_presentor(presentation, presenter):
    slide = presentation.slides[0]
    shape = slide.shapes[3]
    shape.text_frame.clear()

    shape.text_frame.paragraphs[0].add_run()
    shape.text_frame.paragraphs[0].runs[0].text = f'提示人：{presenter}'
    shape.text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    shape.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    shape.text_frame.paragraphs[0].runs[0].font.bold = True
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
    shape.text_frame.paragraphs[0].runs[0].font.brightness = -0.35


def cal_num_of_past_missions_slide(missions, one_page_mission_count):
    num_of_slide = int(len(missions) / one_page_mission_count)
    if (len(missions) % one_page_mission_count) > 0:
        num_of_slide += 1

    if num_of_slide == 0:
        num_of_slide = 1

    return num_of_slide


def setup_past_missions(presentation, missions, host, presenter):
    mission_length_of_each_page = [3] * int(len(missions) / 3)
    if (len(missions) % 3) > 0:
        mission_length_of_each_page.append(len(missions) % 3)

    for slide_idx, mission_length in enumerate(mission_length_of_each_page):
        if slide_idx != 0:
            duplicate_slide(presentation, 2)
            move_slide(presentation, 14+slide_idx, 2)

    for slide_idx, mission_length in enumerate(mission_length_of_each_page):
        slide = presentation.slides[2+slide_idx]

        rows = 5 + mission_length
        cols = 7
        left = Inches(0.06)
        top = Inches(0.61)
        width = Inches(10)
        height = Inches(5)

        shapes = slide.shapes
        table = shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.5)
        table.columns[1].width = Inches(0.9)
        table.columns[2].width = Inches(0.4)
        table.columns[3].width = Inches(1.9)
        table.columns[4].width = Inches(1.1)
        table.columns[5].width = Inches(2.35)
        table.columns[6].width = Inches(2.75)

        table.rows[0].height = Inches(0.4)
        table.rows[1].height = Inches(0.4)
        table.rows[2].height = Inches(0.37)
        table.rows[3].height = Inches(0.37)
        table.rows[4].height = Inches(0.43)

        # cell 0, 0
        table.cell(0, 0).text_frame.clear()
        table.cell(0, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(0, 0).text_frame.paragraphs[0].add_run()
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].text = '內政部空勤務總隊勤務'
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(20)
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.brightness = 1

        table.cell(0, 0).text_frame.add_paragraph()
        table.cell(0, 0).text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
        table.cell(0, 0).text_frame.paragraphs[1].add_run()
        table.cell(0, 0).text_frame.paragraphs[1].runs[0].text = '第三大隊第二隊 '
        table.cell(0, 0).text_frame.paragraphs[1].runs[0].font.name = 'DFKai-SB'
        table.cell(0, 0).text_frame.paragraphs[1].runs[0].font.size = Pt(20)
        table.cell(0, 0).text_frame.paragraphs[1].runs[0].font.bold = True
        table.cell(0, 0).text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        table.cell(0, 0).text_frame.paragraphs[1].runs[0].font.brightness = 1

        table.cell(0, 0).text_frame.paragraphs[1].add_run()
        table.cell(0, 0).text_frame.paragraphs[1].runs[1].text = '每日任務提示紀錄'
        table.cell(0, 0).text_frame.paragraphs[1].runs[1].font.name = 'DFKai-SB'
        table.cell(0, 0).text_frame.paragraphs[1].runs[1].font.size = Pt(20)
        table.cell(0, 0).text_frame.paragraphs[1].runs[1].font.bold = True
        table.cell(0, 0).text_frame.paragraphs[1].runs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        table.cell(0, 0).text_frame.paragraphs[1].runs[1].font.brightness = 1
        mergeCells(table=table, start_row_idx=0, end_row_idx=1, start_col_idx=0, end_col_idx=6)

        # cell 2, 0
        table.cell(2, 0).text_frame.clear()
        table.cell(2, 0).fill.solid()
        table.cell(2, 0).fill.fore_color.rgb = RGBColor(0xB9, 0xCD, 0xE5)
        table.cell(2, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(2, 0).text_frame.paragraphs[0].add_run()
        table.cell(2, 0).text_frame.paragraphs[0].runs[0].text = '主持人'
        table.cell(2, 0).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(2, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        table.cell(2, 0).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(2, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1C, 0x60, 0xB2)
        table.cell(2, 0).text_frame.paragraphs[0].runs[0].font.brightness = 0.3
        mergeCellsHorizontally(table=table, row_idx=2, start_col_idx=0, end_col_idx=2)

        # cell 2, 3
        table.cell(2, 3).text_frame.clear()
        table.cell(2, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(2, 3).text_frame.paragraphs[0].add_run()
        table.cell(2, 3).text_frame.paragraphs[0].runs[0].text = f'{host}'
        table.cell(2, 3).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(2, 3).text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        table.cell(2, 3).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(2, 3).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1C, 0x60, 0xB2)
        table.cell(2, 3).text_frame.paragraphs[0].runs[0].font.brightness = 0.3
        mergeCellsHorizontally(table=table, row_idx=2, start_col_idx=3, end_col_idx=4)

        # cell 2, 5
        table.cell(2, 5).text_frame.clear()
        table.cell(2, 5).fill.solid()
        table.cell(2, 5).fill.fore_color.rgb = RGBColor(0xB9, 0xCD, 0xE5)
        table.cell(2, 5).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(2, 5).text_frame.paragraphs[0].add_run()
        table.cell(2, 5).text_frame.paragraphs[0].runs[0].text = '提示人'
        table.cell(2, 5).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(2, 5).text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        table.cell(2, 5).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(2, 5).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1C, 0x60, 0xB2)
        table.cell(2, 5).text_frame.paragraphs[0].runs[0].font.brightness = 0.3

        # cell 2, 6
        table.cell(2, 6).text_frame.clear()
        table.cell(2, 6).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(2, 6).text_frame.paragraphs[0].add_run()
        table.cell(2, 6).text_frame.paragraphs[0].runs[0].text = f'{presenter}'
        table.cell(2, 6).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(2, 6).text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        table.cell(2, 6).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(2, 6).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1C, 0x60, 0xB2)
        table.cell(2, 6).text_frame.paragraphs[0].runs[0].font.brightness = 0.3

        # cell 3, 0
        table.cell(3, 0).text_frame.clear()
        table.cell(3, 0).fill.background()
        table.cell(3, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        table.cell(3, 0).text_frame.paragraphs[0].add_run()
        table.cell(3, 0).text_frame.paragraphs[0].runs[0].text = '一、前日任務檢討：'
        table.cell(3, 0).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(3, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        table.cell(3, 0).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(3, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1C, 0x60, 0xB2)
        table.cell(3, 0).text_frame.paragraphs[0].runs[0].font.brightness = 0.3
        mergeCellsHorizontally(table=table, row_idx=3, start_col_idx=0, end_col_idx=6)

        # cell 4, 0
        table.cell(4, 0).text_frame.clear()
        table.cell(4, 0).fill.background()
        table.cell(4, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(4, 0).text_frame.paragraphs[0].add_run()
        table.cell(4, 0).text_frame.paragraphs[0].runs[0].text = 'No.'
        table.cell(4, 0).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(4, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(4, 0).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(4, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1C, 0x60, 0xB2)
        table.cell(4, 0).text_frame.paragraphs[0].runs[0].font.brightness = 0.3

        # cell 4, 1
        table.cell(4, 1).text_frame.clear()
        table.cell(4, 1).fill.background()
        table.cell(4, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(4, 1).text_frame.paragraphs[0].add_run()
        table.cell(4, 1).text_frame.paragraphs[0].runs[0].text = '機號'
        table.cell(4, 1).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(4, 1).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(4, 1).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(4, 1).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1C, 0x60, 0xB2)
        table.cell(4, 1).text_frame.paragraphs[0].runs[0].font.brightness = 0.3

        # cell 4, 2
        table.cell(4, 2).text_frame.clear()
        table.cell(4, 2).fill.background()
        table.cell(4, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(4, 2).text_frame.paragraphs[0].add_run()
        table.cell(4, 2).text_frame.paragraphs[0].runs[0].text = '任務項目'
        table.cell(4, 2).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(4, 2).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(4, 2).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(4, 2).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1C, 0x60, 0xB2)
        table.cell(4, 2).text_frame.paragraphs[0].runs[0].font.brightness = 0.3
        mergeCellsHorizontally(table=table, row_idx=4, start_col_idx=2, end_col_idx=3)

        # cell 4, 4
        table.cell(4, 4).text_frame.clear()
        table.cell(4, 4).fill.background()
        table.cell(4, 4).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(4, 4).text_frame.paragraphs[0].add_run()
        table.cell(4, 4).text_frame.paragraphs[0].runs[0].text = '機組人員'
        table.cell(4, 4).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(4, 4).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(4, 4).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(4, 4).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1C, 0x60, 0xB2)
        table.cell(4, 4).text_frame.paragraphs[0].runs[0].font.brightness = 0.3
        mergeCellsHorizontally(table=table, row_idx=4, start_col_idx=4, end_col_idx=5)

        # cell 4, 6
        table.cell(4, 6).text_frame.clear()
        table.cell(4, 6).fill.background()
        table.cell(4, 6).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(4, 6).text_frame.paragraphs[0].add_run()
        table.cell(4, 6).text_frame.paragraphs[0].runs[0].text = '任務時間'
        table.cell(4, 6).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(4, 6).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(4, 6).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(4, 6).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1C, 0x60, 0xB2)
        table.cell(4, 6).text_frame.paragraphs[0].runs[0].font.brightness = 0.3

        for idx in range(mission_length):
            mission = missions[slide_idx*3 + idx]

            # col 0
            table.cell(5+idx, 0).text_frame.clear()
            table.cell(5+idx, 0).fill.background()
            table.cell(5+idx, 0).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(5+idx, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(5+idx, 0).text_frame.paragraphs[0].add_run()
            table.cell(5+idx, 0).text_frame.paragraphs[0].runs[0].text = f'{idx+1}'
            table.cell(5+idx, 0).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
            table.cell(5+idx, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
            table.cell(5+idx, 0).text_frame.paragraphs[0].runs[0].font.bold = True
            table.cell(5+idx, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            table.cell(5+idx, 0).text_frame.paragraphs[0].runs[0].font.brightness = 0

            # col 1
            table.cell(5+idx, 1).text_frame.clear()
            table.cell(5+idx, 1).fill.background()
            table.cell(5+idx, 1).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(5+idx, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(5+idx, 1).text_frame.paragraphs[0].add_run()
            table.cell(5+idx, 1).text_frame.paragraphs[0].runs[0].text = mission['plane-num']
            table.cell(5+idx, 1).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
            table.cell(5+idx, 1).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
            table.cell(5+idx, 1).text_frame.paragraphs[0].runs[0].font.bold = True
            table.cell(5+idx, 1).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            table.cell(5+idx, 1).text_frame.paragraphs[0].runs[0].font.brightness = 0

            # col 2
            table.cell(5+idx, 2).text_frame.clear()
            table.cell(5+idx, 2).fill.background()
            table.cell(5+idx, 2).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(5+idx, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(5+idx, 2).text_frame.paragraphs[0].add_run()
            table.cell(5+idx, 2).text_frame.paragraphs[0].runs[0].text = mission['zh_type']
            table.cell(5+idx, 2).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
            table.cell(5+idx, 2).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
            table.cell(5+idx, 2).text_frame.paragraphs[0].runs[0].font.bold = True
            table.cell(5+idx, 2).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            table.cell(5+idx, 2).text_frame.paragraphs[0].runs[0].font.brightness = 0
            mergeCellsHorizontally(table=table, row_idx=5+idx, start_col_idx=2, end_col_idx=3)

            # col 4
            table.cell(5+idx, 4).text_frame.clear()
            table.cell(5+idx, 4).fill.background()
            table.cell(5+idx, 4).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(5+idx, 4).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(5+idx, 4).text_frame.paragraphs[0].add_run()
            table.cell(5+idx, 4).text_frame.paragraphs[0].runs[0].text = mission['people']
            table.cell(5+idx, 4).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
            table.cell(5+idx, 4).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
            table.cell(5+idx, 4).text_frame.paragraphs[0].runs[0].font.bold = True
            table.cell(5+idx, 4).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            table.cell(5+idx, 4).text_frame.paragraphs[0].runs[0].font.brightness = 0
            mergeCellsHorizontally(table=table, row_idx=5+idx, start_col_idx=4, end_col_idx=5)

            # col 6
            table.cell(5+idx, 6).text_frame.clear()
            table.cell(5+idx, 6).fill.background()
            table.cell(5+idx, 6).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(5+idx, 6).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(5+idx, 6).text_frame.paragraphs[0].add_run()
            table.cell(5+idx, 6).text_frame.paragraphs[0].runs[0].text = f'{mission["time"]}\n{mission["note"]}'
            table.cell(5+idx, 6).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
            table.cell(5+idx, 6).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
            table.cell(5+idx, 6).text_frame.paragraphs[0].runs[0].font.bold = True
            table.cell(5+idx, 6).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            table.cell(5+idx, 6).text_frame.paragraphs[0].runs[0].font.brightness = 0


def setup_today_missions(presentation, start_slide_idx, missions):
    mission_length_of_each_page = [6] * int(len(missions) / 6)
    if (len(missions) % 6) > 0:
        mission_length_of_each_page.append(len(missions) % 6)

    for slide_idx, mission_length in enumerate(mission_length_of_each_page):
        if slide_idx != 0:
            duplicate_slide(presentation, start_slide_idx)
            move_slide(presentation, 11+start_slide_idx+slide_idx, start_slide_idx)

    for slide_idx, mission_length in enumerate(mission_length_of_each_page):
        slide = presentation.slides[start_slide_idx+slide_idx]

        rows = 1 + mission_length
        cols = 4
        left = Inches(0.23)
        top = Inches(0.61)
        width = Inches(9.6)
        height = Inches(5)

        shapes = slide.shapes
        table = shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.87)
        table.columns[1].width = Inches(1.97)
        table.columns[2].width = Inches(4.17)
        table.columns[3].width = Inches(2.53)

        table.rows[0].height = Inches(0.41)

        # cell 0, 0
        table.cell(0, 0).text_frame.clear()
        table.cell(0, 0).fill.background()
        table.cell(0, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(0, 0).text_frame.paragraphs[0].add_run()
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].text = '機號'
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
        table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.brightness = 0.35

        # cell 0, 1
        table.cell(0, 1).text_frame.clear()
        table.cell(0, 1).fill.background()
        table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(0, 1).text_frame.paragraphs[0].add_run()
        table.cell(0, 1).text_frame.paragraphs[0].runs[0].text = '任務項目'
        table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
        table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.brightness = 0.35

        # cell 0, 2
        table.cell(0, 2).text_frame.clear()
        table.cell(0, 2).fill.background()
        table.cell(0, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(0, 2).text_frame.paragraphs[0].add_run()
        table.cell(0, 2).text_frame.paragraphs[0].runs[0].text = '機組人員'
        table.cell(0, 2).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(0, 2).text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        table.cell(0, 2).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(0, 2).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
        table.cell(0, 2).text_frame.paragraphs[0].runs[0].font.brightness = 0.35

        # cell 0, 3
        table.cell(0, 3).text_frame.clear()
        table.cell(0, 3).fill.background()
        table.cell(0, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(0, 3).text_frame.paragraphs[0].add_run()
        table.cell(0, 3).text_frame.paragraphs[0].runs[0].text = '任務時間'
        table.cell(0, 3).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(0, 3).text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        table.cell(0, 3).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(0, 3).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
        table.cell(0, 3).text_frame.paragraphs[0].runs[0].font.brightness = 0.35

        for idx in range(mission_length):
            mission = missions[slide_idx*6 + idx]

            # col 0
            table.cell(1+idx, 0).text_frame.clear()
            table.cell(1+idx, 0).fill.background()
            table.cell(1+idx, 0).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(1+idx, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(1+idx, 0).text_frame.paragraphs[0].add_run()
            table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].text = mission['plane-num']
            table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
            table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(14)
            table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].font.bold = True
            table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].font.brightness = 0

            # col 1
            table.cell(1+idx, 1).text_frame.clear()
            table.cell(1+idx, 1).fill.background()
            table.cell(1+idx, 1).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(1+idx, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(1+idx, 1).text_frame.paragraphs[0].add_run()
            table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].text = mission['zh_type']
            table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
            table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].font.size = Pt(14)
            table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].font.bold = True
            table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].font.brightness = 0

            # col 2
            table.cell(1+idx, 2).text_frame.clear()
            table.cell(1+idx, 2).fill.background()
            table.cell(1+idx, 2).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(1+idx, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(1+idx, 2).text_frame.paragraphs[0].add_run()
            table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].text = mission['people']
            table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
            table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].font.size = Pt(14)
            table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].font.bold = True
            table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].font.brightness = 0

            # col 3
            table.cell(1+idx, 3).text_frame.clear()
            table.cell(1+idx, 3).fill.background()
            table.cell(1+idx, 3).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(1+idx, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(1+idx, 3).text_frame.paragraphs[0].add_run()
            table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].text = f'{mission["time"]}\n{mission["note"]}'
            table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
            table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].font.size = Pt(14)
            table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].font.bold = True
            table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].font.brightness = 0


def setup_plane_status(presentation, start_slide_idx, plane_status_list):
    slide = presentation.slides[start_slide_idx]

    rows = 1 + len(plane_status_list)
    cols = 11
    left = Inches(0)
    top = Inches(0.71)
    width = Inches(10)
    height = Inches(5)

    shapes = slide.shapes
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(0.97)
    table.columns[1].width = Inches(0.85)
    table.columns[2].width = Inches(1.22)
    table.columns[3].width = Inches(0.9)
    table.columns[4].width = Inches(0.74)
    table.columns[5].width = Inches(1.16)
    table.columns[6].width = Inches(1.08)
    table.columns[7].width = Inches(1.01)
    table.columns[8].width = Inches(0.58)
    table.columns[9].width = Inches(0.68)
    table.columns[10].width = Inches(0.8)

    table.rows[0].height = Inches(1.17)

    # cell 0, 0
    table.cell(0, 0).text_frame.clear()
    table.cell(0, 0).fill.background()
    table.cell(0, 0).vertical_anchor = MSO_ANCHOR.MIDDLE
    table.cell(0, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 0).text_frame.paragraphs[0].add_run()
    table.cell(0, 0).text_frame.paragraphs[0].runs[0].text = '機號'
    table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold = True
    table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.brightness = 0

    # cell 0, 1
    table.cell(0, 1).text_frame.clear()
    table.cell(0, 1).fill.background()
    table.cell(0, 1).vertical_anchor = MSO_ANCHOR.MIDDLE
    table.cell(0, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 1).text_frame.paragraphs[0].add_run()
    table.cell(0, 1).text_frame.paragraphs[0].runs[0].text = '狀況'
    table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.bold = True
    table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.brightness = 0

    # cell 0, 2
    table.cell(0, 2).text_frame.clear()
    table.cell(0, 2).fill.background()
    table.cell(0, 2).vertical_anchor = MSO_ANCHOR.MIDDLE
    table.cell(0, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 2).text_frame.paragraphs[0].add_run()
    table.cell(0, 2).text_frame.paragraphs[0].runs[0].text = '檢修日期'
    table.cell(0, 2).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    table.cell(0, 2).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    table.cell(0, 2).text_frame.paragraphs[0].runs[0].font.bold = True
    table.cell(0, 2).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    table.cell(0, 2).text_frame.paragraphs[0].runs[0].font.brightness = 0

    # cell 0, 3
    table.cell(0, 3).text_frame.clear()
    table.cell(0, 3).fill.background()
    table.cell(0, 3).vertical_anchor = MSO_ANCHOR.MIDDLE
    table.cell(0, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 3).text_frame.paragraphs[0].add_run()
    table.cell(0, 3).text_frame.paragraphs[0].runs[0].text = '位置'
    table.cell(0, 3).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    table.cell(0, 3).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    table.cell(0, 3).text_frame.paragraphs[0].runs[0].font.bold = True
    table.cell(0, 3).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    table.cell(0, 3).text_frame.paragraphs[0].runs[0].font.brightness = 0

    # cell 0, 4
    table.cell(0, 4).text_frame.clear()
    table.cell(0, 4).fill.background()
    table.cell(0, 4).vertical_anchor = MSO_ANCHOR.MIDDLE
    table.cell(0, 4).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 4).text_frame.paragraphs[0].add_run()
    table.cell(0, 4).text_frame.paragraphs[0].runs[0].text = '前日飛行時間'
    table.cell(0, 4).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    table.cell(0, 4).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    table.cell(0, 4).text_frame.paragraphs[0].runs[0].font.bold = True
    table.cell(0, 4).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    table.cell(0, 4).text_frame.paragraphs[0].runs[0].font.brightness = 0

    # cell 0, 5
    table.cell(0, 5).text_frame.clear()
    table.cell(0, 5).fill.background()
    table.cell(0, 5).vertical_anchor = MSO_ANCHOR.MIDDLE
    table.cell(0, 5).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 5).text_frame.paragraphs[0].add_run()
    table.cell(0, 5).text_frame.paragraphs[0].runs[0].text = '飛機時間'
    table.cell(0, 5).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    table.cell(0, 5).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    table.cell(0, 5).text_frame.paragraphs[0].runs[0].font.bold = True
    table.cell(0, 5).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    table.cell(0, 5).text_frame.paragraphs[0].runs[0].font.brightness = 0

    # cell 0, 6
    table.cell(0, 6).text_frame.clear()
    table.cell(0, 6).fill.background()
    table.cell(0, 6).vertical_anchor = MSO_ANCHOR.MIDDLE
    table.cell(0, 6).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 6).text_frame.paragraphs[0].add_run()
    table.cell(0, 6).text_frame.paragraphs[0].runs[0].text = '發動機時間'
    table.cell(0, 6).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    table.cell(0, 6).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    table.cell(0, 6).text_frame.paragraphs[0].runs[0].font.bold = True
    table.cell(0, 6).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    table.cell(0, 6).text_frame.paragraphs[0].runs[0].font.brightness = 0

    # cell 0, 7
    table.cell(0, 7).text_frame.clear()
    table.cell(0, 7).fill.background()
    table.cell(0, 7).vertical_anchor = MSO_ANCHOR.MIDDLE
    table.cell(0, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 7).text_frame.paragraphs[0].add_run()
    table.cell(0, 7).text_frame.paragraphs[0].runs[0].text = '距週(階)檢時間'
    table.cell(0, 7).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    table.cell(0, 7).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    table.cell(0, 7).text_frame.paragraphs[0].runs[0].font.bold = True
    table.cell(0, 7).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    table.cell(0, 7).text_frame.paragraphs[0].runs[0].font.brightness = 0

    # cell 0, 8
    table.cell(0, 8).text_frame.clear()
    table.cell(0, 8).fill.background()
    table.cell(0, 8).vertical_anchor = MSO_ANCHOR.MIDDLE
    table.cell(0, 8).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 8).text_frame.paragraphs[0].add_run()
    table.cell(0, 8).text_frame.paragraphs[0].runs[0].text = '人員吊掛'
    table.cell(0, 8).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    table.cell(0, 8).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    table.cell(0, 8).text_frame.paragraphs[0].runs[0].font.bold = True
    table.cell(0, 8).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    table.cell(0, 8).text_frame.paragraphs[0].runs[0].font.brightness = 0

    # cell 0, 9
    table.cell(0, 9).text_frame.clear()
    table.cell(0, 9).fill.background()
    table.cell(0, 9).vertical_anchor = MSO_ANCHOR.MIDDLE
    table.cell(0, 9).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 9).text_frame.paragraphs[0].add_run()
    table.cell(0, 9).text_frame.paragraphs[0].runs[0].text = '緊急浮筒'
    table.cell(0, 9).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    table.cell(0, 9).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    table.cell(0, 9).text_frame.paragraphs[0].runs[0].font.bold = True
    table.cell(0, 9).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    table.cell(0, 9).text_frame.paragraphs[0].runs[0].font.brightness = 0

    # cell 0, 10
    table.cell(0, 10).text_frame.clear()
    table.cell(0, 10).fill.background()
    table.cell(0, 10).vertical_anchor = MSO_ANCHOR.MIDDLE
    table.cell(0, 10).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(0, 10).text_frame.paragraphs[0].add_run()
    table.cell(0, 10).text_frame.paragraphs[0].runs[0].text = '油量'
    table.cell(0, 10).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
    table.cell(0, 10).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    table.cell(0, 10).text_frame.paragraphs[0].runs[0].font.bold = True
    table.cell(0, 10).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    table.cell(0, 10).text_frame.paragraphs[0].runs[0].font.brightness = 0

    for idx, plane_status in enumerate(plane_status_list):
        # col 0
        table.cell(1+idx, 0).text_frame.clear()
        table.cell(1+idx, 0).fill.background()
        table.cell(1+idx, 0).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(1+idx, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(1+idx, 0).text_frame.paragraphs[0].add_run()
        table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].text = plane_status['plane-num']
        table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        table.cell(1+idx, 0).text_frame.paragraphs[0].runs[0].font.brightness = 0

        # col 1
        table.cell(1+idx, 1).text_frame.clear()
        table.cell(1+idx, 1).fill.background()
        table.cell(1+idx, 1).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(1+idx, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(1+idx, 1).text_frame.paragraphs[0].add_run()
        table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].text = plane_status['status']
        table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        table.cell(1+idx, 1).text_frame.paragraphs[0].runs[0].font.brightness = 0

        # col 2
        table.cell(1+idx, 2).text_frame.clear()
        table.cell(1+idx, 2).fill.background()
        table.cell(1+idx, 2).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(1+idx, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(1+idx, 2).text_frame.paragraphs[0].add_run()
        table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].text = plane_status['check_date']
        table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        table.cell(1+idx, 2).text_frame.paragraphs[0].runs[0].font.brightness = 0

        # col 3
        table.cell(1+idx, 3).text_frame.clear()
        table.cell(1+idx, 3).fill.background()
        table.cell(1+idx, 3).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(1+idx, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(1+idx, 3).text_frame.paragraphs[0].add_run()
        table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].text = plane_status['position']
        table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        table.cell(1+idx, 3).text_frame.paragraphs[0].runs[0].font.brightness = 0

        # col 4
        table.cell(1+idx, 4).text_frame.clear()
        table.cell(1+idx, 4).fill.background()
        table.cell(1+idx, 4).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(1+idx, 4).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(1+idx, 4).text_frame.paragraphs[0].add_run()
        table.cell(1+idx, 4).text_frame.paragraphs[0].runs[0].text = plane_status['yesterday_time']
        table.cell(1+idx, 4).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(1+idx, 4).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(1+idx, 4).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(1+idx, 4).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        table.cell(1+idx, 4).text_frame.paragraphs[0].runs[0].font.brightness = 0

        # col 5
        table.cell(1+idx, 5).text_frame.clear()
        table.cell(1+idx, 5).fill.background()
        table.cell(1+idx, 5).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(1+idx, 5).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(1+idx, 5).text_frame.paragraphs[0].add_run()
        table.cell(1+idx, 5).text_frame.paragraphs[0].runs[0].text = plane_status['plane_time']
        table.cell(1+idx, 5).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(1+idx, 5).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(1+idx, 5).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(1+idx, 5).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        table.cell(1+idx, 5).text_frame.paragraphs[0].runs[0].font.brightness = 0

        # col 6
        table.cell(1+idx, 6).text_frame.clear()
        table.cell(1+idx, 6).fill.background()
        table.cell(1+idx, 6).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(1+idx, 6).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(1+idx, 6).text_frame.paragraphs[0].add_run()
        table.cell(1+idx, 6).text_frame.paragraphs[0].runs[0].text = plane_status['engine_time']
        table.cell(1+idx, 6).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(1+idx, 6).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(1+idx, 6).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(1+idx, 6).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        table.cell(1+idx, 6).text_frame.paragraphs[0].runs[0].font.brightness = 0

        # col 7
        table.cell(1+idx, 7).text_frame.clear()
        table.cell(1+idx, 7).fill.background()
        table.cell(1+idx, 7).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(1+idx, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(1+idx, 7).text_frame.paragraphs[0].add_run()
        table.cell(1+idx, 7).text_frame.paragraphs[0].runs[0].text = plane_status['distance_check_time']
        table.cell(1+idx, 7).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(1+idx, 7).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(1+idx, 7).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(1+idx, 7).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        table.cell(1+idx, 7).text_frame.paragraphs[0].runs[0].font.brightness = 0

        # col 8
        table.cell(1+idx, 8).text_frame.clear()
        table.cell(1+idx, 8).fill.background()
        table.cell(1+idx, 8).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(1+idx, 8).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(1+idx, 8).text_frame.paragraphs[0].add_run()
        table.cell(1+idx, 8).text_frame.paragraphs[0].runs[0].text = plane_status['person_hang']
        table.cell(1+idx, 8).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(1+idx, 8).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(1+idx, 8).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(1+idx, 8).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        table.cell(1+idx, 8).text_frame.paragraphs[0].runs[0].font.brightness = 0

        # col 9
        table.cell(1+idx, 9).text_frame.clear()
        table.cell(1+idx, 9).fill.background()
        table.cell(1+idx, 9).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(1+idx, 9).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(1+idx, 9).text_frame.paragraphs[0].add_run()
        table.cell(1+idx, 9).text_frame.paragraphs[0].runs[0].text = plane_status['emergency_buoy']
        table.cell(1+idx, 9).text_frame.paragraphs[0].runs[0].font.name = 'DFKai-SB'
        table.cell(1+idx, 9).text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        table.cell(1+idx, 9).text_frame.paragraphs[0].runs[0].font.bold = True
        table.cell(1+idx, 9).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        table.cell(1+idx, 9).text_frame.paragraphs[0].runs[0].font.brightness = 0

        # col 10
        table.cell(1+idx, 10).text_frame.clear()
        table.cell(1+idx, 10).fill.background()
        table.cell(1+idx, 10).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(1+idx, 10).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def setup_weather_overview(presentation, start_slide_idx):
    req = requests.get('http://www.cwb.gov.tw/V7/forecast/txt/w01.htm')
    req.encoding = 'utf-8'
    soup = BeautifulSoup(req.text, 'html.parser')
    paragraph_start_idx = 0
    paragraph_end_idx = 0
    paragraph_list = soup.find('pre').string.strip().split()
    for idx, paragraph in enumerate(paragraph_list):
        if paragraph.startswith('二、'):
            paragraph_start_idx = idx
        elif paragraph.startswith('三、'):
            paragraph_end_idx = idx
    weather_overview = ''.join(paragraph_list[paragraph_start_idx : paragraph_end_idx])
    slide = presentation.slides[start_slide_idx]
    left = Inches(2.4)
    top = Inches(1.47)
    width = Inches(5)
    height = Inches(3)
    textBox = slide.shapes.add_textbox(left, top, width, height)
    textBox.text = weather_overview


def setup_weather_img(presentation, start_slide_idx):
    if os.path.isfile('sateelite.jpg'):
        slide = presentation.slides[start_slide_idx]
        left = Inches(2.24)
        top = Inches(0)
        width = height = Inches(5.63)
        slide.shapes.add_picture('sateelite.jpg', left, top, width, height)

    if os.path.isfile('rader.jpg'):
        slide = presentation.slides[start_slide_idx+1]
        left = Inches(1.93)
        top = Inches(0)
        width = Inches(7.5)
        height = Inches(5.63)
        slide.shapes.add_picture('rader.jpg', left, top, width, height)

    if os.path.isfile('plane_alert.png'):
        slide = presentation.slides[start_slide_idx+4]
        left = Inches(3.37)
        top = Inches(0)
        width = Inches(4.5)
        height = Inches(5.63)
        slide.shapes.add_picture('plane_alert.png', left, top, width, height)


def setup_weather_metar(presentation, start_slide_idx):
    if not os.path.isfile('weather_metar.json'):
        return

    weather_metar_list = json.load(open('weather_metar.json', 'r'))
    weather_metar_json = { metar['station']: metar['message'] for metar in weather_metar_list}
    slide = presentation.slides[start_slide_idx]
    left = Inches(0.01)
    top = Inches(0.37)
    width = Inches(10)
    height = Inches(5)
    textBox = slide.shapes.add_textbox(left, top, width, height)
    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[0].text = 'RCKH'
    textBox.text_frame.paragraphs[0].font.size = Pt(24)
    textBox.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    textBox.text_frame.paragraphs[0].font.brightness = 0.43

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[1].text = weather_metar_json.get('RCKH', '')
    textBox.text_frame.paragraphs[1].font.size = Pt(24)
    textBox.text_frame.paragraphs[1].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    textBox.text_frame.paragraphs[1].font.brightness = 0

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[2].text = 'RCBS'
    textBox.text_frame.paragraphs[2].font.size = Pt(24)
    textBox.text_frame.paragraphs[2].font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    textBox.text_frame.paragraphs[2].font.brightness = 0.43

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[3].text = weather_metar_json.get('RCBS', '')
    textBox.text_frame.paragraphs[3].font.size = Pt(24)
    textBox.text_frame.paragraphs[3].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    textBox.text_frame.paragraphs[3].font.brightness = 0

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[4].text = 'RCDC'
    textBox.text_frame.paragraphs[4].font.size = Pt(24)
    textBox.text_frame.paragraphs[4].font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    textBox.text_frame.paragraphs[4].font.brightness = 0.43

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[5].text = weather_metar_json.get('RCDC', '')
    textBox.text_frame.paragraphs[5].font.size = Pt(24)
    textBox.text_frame.paragraphs[5].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    textBox.text_frame.paragraphs[5].font.brightness = 0

    slide = presentation.slides[start_slide_idx+1]
    left = Inches(0.01)
    top = Inches(0.37)
    width = Inches(10)
    height = Inches(5)
    textBox = slide.shapes.add_textbox(left, top, width, height)
    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[0].text = 'RCFN'
    textBox.text_frame.paragraphs[0].font.size = Pt(24)
    textBox.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    textBox.text_frame.paragraphs[0].font.brightness = 0.43

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[1].text = weather_metar_json.get('RCFN', '')
    textBox.text_frame.paragraphs[1].font.size = Pt(24)
    textBox.text_frame.paragraphs[1].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    textBox.text_frame.paragraphs[1].font.brightness = 0

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[2].text = 'RCLY'
    textBox.text_frame.paragraphs[2].font.size = Pt(24)
    textBox.text_frame.paragraphs[2].font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    textBox.text_frame.paragraphs[2].font.brightness = 0.43

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[3].text = weather_metar_json.get('RCLY', '')
    textBox.text_frame.paragraphs[3].font.size = Pt(24)
    textBox.text_frame.paragraphs[3].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    textBox.text_frame.paragraphs[3].font.brightness = 0

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[4].text = 'RCNN'
    textBox.text_frame.paragraphs[4].font.size = Pt(24)
    textBox.text_frame.paragraphs[4].font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    textBox.text_frame.paragraphs[4].font.brightness = 0.43

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[5].text = weather_metar_json.get('RCNN', '')
    textBox.text_frame.paragraphs[5].font.size = Pt(24)
    textBox.text_frame.paragraphs[5].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    textBox.text_frame.paragraphs[5].font.brightness = 0

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[6].text = 'RCQC'
    textBox.text_frame.paragraphs[6].font.size = Pt(24)
    textBox.text_frame.paragraphs[6].font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    textBox.text_frame.paragraphs[6].font.brightness = 0.43

    textBox.text_frame.add_paragraph()
    textBox.text_frame.paragraphs[6].text = weather_metar_json.get('RCQC', '')
    textBox.text_frame.paragraphs[6].font.size = Pt(24)
    textBox.text_frame.paragraphs[6].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    textBox.text_frame.paragraphs[6].font.brightness = 0


if __name__ == '__main__':
    input_encoding = chardet.detect(open('pptx_input.json', 'rb').read())['encoding']
    input_json = json.load(open('pptx_input.json', 'r', encoding=input_encoding))

    past_missions = input_json['past_missions']
    today_missions = input_json['today_missions']
    plane_status_list = input_json['plane_status_list']

    prs = Presentation('./template.pptx')
    setup_past_missions(prs, past_missions, 'somebody', 'somebody')

    start_today_slide_idx = 2 + cal_num_of_past_missions_slide(past_missions, 3)
    setup_today_missions(prs, start_today_slide_idx, today_missions)

    start_plane_status_slide_idx = start_today_slide_idx + cal_num_of_past_missions_slide(today_missions, 6)
    setup_plane_status(prs, start_plane_status_slide_idx, plane_status_list)

    start_weather_overview_slide_idx = start_plane_status_slide_idx + 1
    setup_weather_overview(prs, start_weather_overview_slide_idx)

    start_weather_img_slide_idx = start_weather_overview_slide_idx + 1
    setup_weather_img(prs, start_weather_img_slide_idx)

    start_weather_metar_slide_idx = start_weather_img_slide_idx + 2
    setup_weather_metar(prs, start_weather_metar_slide_idx)
    prs.save('output.pptx')
